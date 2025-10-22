import io
import json
import base64
import re
import datetime
from collections import Counter
import heapq
import tempfile
import os

# Flask for web server
from flask import Flask, request, jsonify
from flask_cors import CORS

# Libraries for analysis
import PyPDF2
from pptx import Presentation
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords
from nltk.tag import pos_tag
from nltk.sentiment.vader import SentimentIntensityAnalyzer
import textstat
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from sklearn.cluster import KMeans
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from wordcloud import WordCloud
import matplotlib.pyplot as plt
import networkx as nx
from textblob import TextBlob
import numpy as np
from textstat import flesch_reading_ease, flesch_kincaid_grade, smog_index, coleman_liau_index
import language_tool_python

# --- NLTK Data Download ---
try:
    nltk.data.find("tokenizers/punkt")
except LookupError:
    nltk.download("punkt")
try:
    nltk.data.find("corpora/stopwords")
except LookupError:
    nltk.download("stopwords")
try:
    nltk.data.find("taggers/averaged_perceptron_tagger")
except LookupError:
    nltk.download("averaged_perceptron_tagger")
    
# Additional download for the specific resource
try:
    nltk.data.find("taggers/averaged_perceptron_tagger_eng")
except LookupError:
    nltk.download("averaged_perceptron_tagger")
try:
    nltk.data.find("sentiment/vader_lexicon")
except LookupError:
    nltk.download("vader_lexicon")

plt.switch_backend('Agg')

# --- Flask App Initialization ---
app = Flask(__name__)
CORS(app)

# --- Enhanced Analysis Logic ---
class EnhancedDocuMentorProAnalyzer:
    """Enhanced backend text processing and analysis helper."""
    
    def __init__(self):
        self.stop_words = set(stopwords.words('english'))
        self.sia = SentimentIntensityAnalyzer()
        self.grammar_tool = language_tool_python.LanguageTool('en-US')
        
        # Enhanced stopwords for better keyword extraction
        self.enhanced_stopwords = self.stop_words.union({
            'however', 'therefore', 'moreover', 'furthermore', 'nevertheless',
            'consequently', 'additionally', 'similarly', 'specifically'
        })

    def extract_text(self, file):
        text = ""
        filename = file.filename.lower()
        file_size = len(file.read())
        file.seek(0)  # Reset file pointer
        
        try:
            if filename.endswith(".pdf"):
                pdf_reader = PyPDF2.PdfReader(file)
                num_pages = len(pdf_reader.pages)
                for page in pdf_reader.pages:
                    text += (page.extract_text() or "")
                metadata = {
                    "pages": num_pages,
                    "file_size": file_size,
                    "file_type": "PDF"
                }
            elif filename.endswith(".pptx"):
                prs = Presentation(file)
                num_slides = len(prs.slides)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text += shape.text + "\n"
                metadata = {
                    "pages": num_slides,
                    "file_size": file_size,
                    "file_type": "PPTX"
                }
            elif filename.endswith(".txt"):
                raw = file.read()
                text = raw.decode('utf-8', errors='replace')
                metadata = {
                    "pages": 1,
                    "file_size": file_size,
                    "file_type": "TXT"
                }
            else:
                return None, {"error": "Unsupported file format"}
                
            return text.strip(), metadata
        except Exception as e:
            print(f"Error extracting text: {e}")
            return None, {"error": str(e)}

    def get_enhanced_stats(self, text):
        """Enhanced statistics with more metrics"""
        if not text or len(text.strip()) == 0:
            return self._get_empty_stats()
            
        words = word_tokenize(text)
        sentences = sent_tokenize(text)
        paragraphs = [p for p in text.split('\n\n') if p.strip()]
        
        # Character and word analysis
        char_count = len(text)
        word_count = len(words)
        sentence_count = len(sentences)
        paragraph_count = len(paragraphs)
        
        # Readability scores
        try:
            flesch_ease = flesch_reading_ease(text)
            flesch_grade = flesch_kincaid_grade(text)
            smog = smog_index(text)
            coleman_liau = coleman_liau_index(text)
            avg_sentence_length = word_count / sentence_count if sentence_count > 0 else 0
            avg_word_length = char_count / word_count if word_count > 0 else 0
        except Exception:
            flesch_ease = flesch_grade = smog = coleman_liau = 0
            avg_sentence_length = avg_word_length = 0

        # Vocabulary analysis
        unique_words = set(word.lower() for word in words if word.isalpha())
        lexical_diversity = len(unique_words) / word_count if word_count > 0 else 0
        
        return {
            "basic": {
                "word_count": word_count,
                "sentence_count": sentence_count,
                "paragraph_count": paragraph_count,
                "character_count": char_count,
                "avg_sentence_length": round(avg_sentence_length, 1),
                "avg_word_length": round(avg_word_length, 1),
            },
            "readability": {
                "flesch_reading_ease": round(flesch_ease, 1),
                "flesch_kincaid_grade": round(flesch_grade, 1),
                "smog_index": round(smog, 1),
                "coleman_liau_index": round(coleman_liau, 1),
            },
            "vocabulary": {
                "unique_words": len(unique_words),
                "lexical_diversity": round(lexical_diversity, 3),
            }
        }

    def _get_empty_stats(self):
        return {
            "basic": {"word_count": 0, "sentence_count": 0, "paragraph_count": 0, 
                     "character_count": 0, "avg_sentence_length": 0, "avg_word_length": 0},
            "readability": {"flesch_reading_ease": 0, "flesch_kincaid_grade": 0, 
                           "smog_index": 0, "coleman_liau_index": 0},
            "vocabulary": {"unique_words": 0, "lexical_diversity": 0}
        }

    def get_enhanced_sentiment(self, text):
        """Enhanced sentiment analysis with emotion detection"""
        if not text:
            return self._get_empty_sentiment()
            
        # Basic sentiment scores
        vader_scores = self.sia.polarity_scores(text)
        blob = TextBlob(text)
        
        # Emotion analysis (simplified)
        positive_words = ['good', 'great', 'excellent', 'amazing', 'wonderful', 'fantastic']
        negative_words = ['bad', 'terrible', 'awful', 'horrible', 'disappointing']
        
        words = word_tokenize(text.lower())
        positive_count = sum(1 for word in words if word in positive_words)
        negative_count = sum(1 for word in words if word in negative_words)
        
        emotion_ratio = positive_count / (positive_count + negative_count) if (positive_count + negative_count) > 0 else 0.5
        
        return {
            "vader_scores": {
                "positive": round(vader_scores['pos'] * 100, 2),
                "neutral": round(vader_scores['neu'] * 100, 2),
                "negative": round(vader_scores['neg'] * 100, 2),
                "compound": round(vader_scores['compound'], 3)
            },
            "textblob_scores": {
                "polarity": round(blob.sentiment.polarity, 3),
                "subjectivity": round(blob.sentiment.subjectivity, 3)
            },
            "emotion_analysis": {
                "positive_words": positive_count,
                "negative_words": negative_count,
                "emotion_ratio": round(emotion_ratio, 3)
            }
        }

    def _get_empty_sentiment(self):
        return {
            "vader_scores": {"positive": 0, "neutral": 100, "negative": 0, "compound": 0},
            "textblob_scores": {"polarity": 0, "subjectivity": 0},
            "emotion_analysis": {"positive_words": 0, "negative_words": 0, "emotion_ratio": 0.5}
        }

    def get_writing_quality(self, text):
        """Analyze writing quality and grammar"""
        if not text:
            return {"grammar_errors": 0, "error_categories": {}, "quality_score": 100}
            
        # Grammar checking
        matches = self.grammar_tool.check(text)
        grammar_errors = len(matches)
        
        # Categorize errors
        error_categories = {}
        for match in matches:
            category = match.category
            error_categories[category] = error_categories.get(category, 0) + 1
        
        # Calculate quality score (0-100)
        word_count = len(word_tokenize(text))
        error_density = grammar_errors / word_count if word_count > 0 else 0
        quality_score = max(0, 100 - (error_density * 1000))  # Scale factor
        
        return {
            "grammar_errors": grammar_errors,
            "error_categories": error_categories,
            "quality_score": round(quality_score, 1),
            "error_density": round(error_density, 4)
        }

    def get_content_structure(self, text):
        """Analyze document structure and organization"""
        if not text:
            return {"sections": 0, "headings": [], "paragraph_distribution": []}
            
        sentences = sent_tokenize(text)
        paragraphs = [p for p in text.split('\n\n') if p.strip()]
        
        # Detect potential headings (sentences that are short and at paragraph start)
        potential_headings = []
        for para in paragraphs:
            first_sentence = sent_tokenize(para)[0] if sent_tokenize(para) else ""
            if len(first_sentence.split()) <= 10 and len(first_sentence) < 100:
                potential_headings.append(first_sentence[:50] + "..." if len(first_sentence) > 50 else first_sentence)
        
        # Paragraph length distribution
        para_lengths = [len(word_tokenize(p)) for p in paragraphs]
        
        return {
            "sections": len(potential_headings) or 1,
            "headings": potential_headings[:10],  # Top 10 potential headings
            "paragraph_distribution": {
                "short_paras": len([p for p in para_lengths if p < 50]),
                "medium_paras": len([p for p in para_lengths if 50 <= p <= 150]),
                "long_paras": len([p for p in para_lengths if p > 150])
            },
            "avg_paragraph_length": round(sum(para_lengths) / len(para_lengths), 1) if para_lengths else 0
        }

    def get_topic_modeling(self, text, num_topics=3):
        """Simple topic modeling using TF-IDF and clustering"""
        if not text:
            return {"topics": [], "topic_keywords": {}}
            
        sentences = sent_tokenize(text)
        if len(sentences) < num_topics:
            num_topics = len(sentences)
            
        # Use TF-IDF to extract important terms
        vectorizer = TfidfVectorizer(max_features=50, stop_words='english')
        try:
            tfidf_matrix = vectorizer.fit_transform(sentences)
            feature_names = vectorizer.get_feature_names_out()
            
            # Simple clustering for topics
            if len(sentences) > num_topics:
                kmeans = KMeans(n_clusters=num_topics, random_state=42)
                kmeans.fit(tfidf_matrix)
                
                # Get top words for each cluster
                topics = {}
                for i in range(num_topics):
                    centroid = kmeans.cluster_centers_[i]
                    top_indices = centroid.argsort()[-5:][::-1]
                    topics[f"Topic {i+1}"] = [feature_names[idx] for idx in top_indices]
            else:
                topics = {"Topic 1": feature_names[:5].tolist()}
                
            return {"topics": list(topics.keys()), "topic_keywords": topics}
        except Exception as e:
            print(f"Topic modeling error: {e}")
            return {"topics": ["General"], "topic_keywords": {"General": ["content", "document", "information"]}}

    def summarize_text(self, text, num_sentences=5):
        """Enhanced text summarization"""
        if not text:
            return ""
        sentences = sent_tokenize(text)
        if len(sentences) <= num_sentences:
            return text
            
        # Use TextRank-like algorithm
        word_frequencies = Counter(w.lower() for w in word_tokenize(text) 
                                 if w.isalpha() and w.lower() not in self.enhanced_stopwords)
        if not word_frequencies:
            return "Could not generate summary."
            
        max_freq = max(word_frequencies.values())
        for word in word_frequencies:
            word_frequencies[word] /= max_freq
            
        sentence_scores = {}
        for sent in sentences:
            for word in word_tokenize(sent.lower()):
                if word in word_frequencies:
                    sentence_scores[sent] = sentence_scores.get(sent, 0) + word_frequencies[word]
            # Normalize by sentence length
            sentence_scores[sent] = sentence_scores.get(sent, 0) / len(word_tokenize(sent))
            
        summary_sentences = heapq.nlargest(num_sentences, sentence_scores, key=sentence_scores.get)
        return ' '.join(summary_sentences)

    def get_enhanced_keywords(self, text, top_n=20):
        """Enhanced keyword extraction with phrases"""
        if not text:
            return []
            
        # Single word keywords
        words = [w.lower() for w in word_tokenize(text) 
                if w.isalpha() and w.lower() not in self.enhanced_stopwords and len(w) > 2]
        word_freq = Counter(words)
        
        # Extract noun phrases
        sentences = sent_tokenize(text)
        noun_phrases = []
        for sent in sentences:
            tokens = word_tokenize(sent)
            pos_tags = pos_tag(tokens)
            
            # Simple noun phrase detection (adjective + noun or noun + noun)
            phrases = []
            current_phrase = []
            for word, pos in pos_tags:
                if pos.startswith('NN') or pos.startswith('JJ'):
                    current_phrase.append(word)
                else:
                    if len(current_phrase) > 1:
                        phrases.append(' '.join(current_phrase))
                    current_phrase = []
            
            if len(current_phrase) > 1:
                phrases.append(' '.join(current_phrase))
                
            noun_phrases.extend(phrases)
        
        # Combine single words and phrases
        all_keywords = [(word, freq) for word, freq in word_freq.most_common(top_n)]
        phrase_freq = Counter(noun_phrases)
        all_keywords.extend([(phrase, freq) for phrase, freq in phrase_freq.most_common(10)])
        
        return all_keywords[:top_n]

    def generate_insights(self, text, stats, sentiment, keywords):
        """Generate AI-powered insights about the document"""
        insights = []
        
        if not text:
            return insights
            
        word_count = stats["basic"]["word_count"]
        readability = stats["readability"]["flesch_reading_ease"]
        polarity = sentiment["textblob_scores"]["polarity"]
        
        # Readability insights
        if readability > 80:
            insights.append("Your document is very easy to read, suitable for a general audience.")
        elif readability > 60:
            insights.append("The document has good readability, appropriate for most readers.")
        elif readability > 30:
            insights.append("Consider simplifying some sentences to improve readability.")
        else:
            insights.append("The document may be difficult for many readers to understand.")
            
        # Length insights
        if word_count < 300:
            insights.append("This is a concise document. Consider adding more detail if needed.")
        elif word_count > 2000:
            insights.append("This is a comprehensive document. Consider breaking it into sections.")
            
        # Sentiment insights
        if polarity > 0.3:
            insights.append("The document has a positive tone, which engages readers well.")
        elif polarity < -0.3:
            insights.append("The negative tone might affect how readers perceive the content.")
        else:
            insights.append("The document maintains a neutral, professional tone.")
            
        # Keyword insights
        if keywords:
            top_keywords = [kw[0] for kw in keywords[:3]]
            insights.append(f"The main topics appear to be: {', '.join(top_keywords)}")
            
        return insights

    def create_enhanced_visualizations(self, text, keywords, stats, sentiment):
        """Create multiple visualization types"""
        visualizations = {}
        
        # Word Cloud
        if keywords:
            try:
                wc = WordCloud(width=800, height=400, background_color="rgba(255, 255, 255, 0)",
                             colormap='viridis', mode="RGBA").generate_from_frequencies(dict(keywords[:30]))
                img_buffer = io.BytesIO()
                wc.to_image().save(img_buffer, format='PNG')
                img_str = base64.b64encode(img_buffer.getvalue()).decode()
                visualizations["wordcloud"] = f"data:image/png;base64,{img_str}"
            except Exception as e:
                print(f"WordCloud error: {e}")
                visualizations["wordcloud"] = None
        
        # Create readability gauge chart
        try:
            readability_score = stats["readability"]["flesch_reading_ease"]
            fig = go.Figure(go.Indicator(
                mode = "gauge+number+delta",
                value = readability_score,
                domain = {'x': [0, 1], 'y': [0, 1]},
                title = {'text': "Readability Score"},
                gauge = {
                    'axis': {'range': [None, 100]},
                    'bar': {'color': "darkblue"},
                    'steps': [
                        {'range': [0, 30], 'color': "red"},
                        {'range': [30, 70], 'color': "yellow"},
                        {'range': [70, 100], 'color': "green"}],
                    'threshold': {
                        'line': {'color': "red", 'width': 4},
                        'thickness': 0.75,
                        'value': 60}}))
            
            fig.update_layout(height=300, paper_bgcolor='rgba(0,0,0,0)')
            visualizations["readability_gauge"] = fig.to_json()
        except Exception as e:
            print(f"Gauge chart error: {e}")
            
        return visualizations

# Instantiate the enhanced analyzer
analyzer = EnhancedDocuMentorProAnalyzer()

# --- Enhanced API Endpoints ---
@app.route('/analyze', methods=['POST'])
def analyze_document():
    """
    Enhanced API endpoint for document analysis
    """
    if 'file' not in request.files:
        return jsonify({"error": "No file part in the request"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No file selected for uploading"}), 400

    try:
        # Extract text and metadata
        text, metadata = analyzer.extract_text(file)
        if not text:
            return jsonify({"error": "Could not extract text from the document."}), 500

        # Perform comprehensive analysis
        stats = analyzer.get_enhanced_stats(text)
        sentiment = analyzer.get_enhanced_sentiment(text)
        writing_quality = analyzer.get_writing_quality(text)
        content_structure = analyzer.get_content_structure(text)
        topics = analyzer.get_topic_modeling(text)
        summary = analyzer.summarize_text(text)
        keywords = analyzer.get_enhanced_keywords(text)
        insights = analyzer.generate_insights(text, stats, sentiment, keywords)
        
        # Generate visualizations
        visualizations = analyzer.create_enhanced_visualizations(text, keywords, stats, sentiment)

        # Assemble comprehensive response
        result = {
            "filename": file.filename,
            "metadata": metadata,
            "stats": stats,
            "sentiment": sentiment,
            "writing_quality": writing_quality,
            "content_structure": content_structure,
            "topics": topics,
            "summary": summary,
            "keywords": keywords,
            "insights": insights,
            "visualizations": visualizations,
            "analysis_timestamp": datetime.datetime.now().isoformat()
        }
        
        return jsonify(result)

    except Exception as e:
        print(f"An error occurred during analysis: {e}")
        return jsonify({"error": "An internal error occurred during analysis."}), 500

@app.route('/batch-analyze', methods=['POST'])
def batch_analyze_documents():
    """
    Analyze multiple documents at once
    """
    if 'files' not in request.files:
        return jsonify({"error": "No files uploaded"}), 400
        
    files = request.files.getlist('files')
    if not files or all(file.filename == '' for file in files):
        return jsonify({"error": "No files selected"}), 400

    results = []
    for file in files:
        if file.filename:
            try:
                text, metadata = analyzer.extract_text(file)
                if text:
                    stats = analyzer.get_enhanced_stats(text)
                    sentiment = analyzer.get_enhanced_sentiment(text)
                    
                    results.append({
                        "filename": file.filename,
                        "word_count": stats["basic"]["word_count"],
                        "readability": stats["readability"]["flesch_reading_ease"],
                        "sentiment": sentiment["textblob_scores"]["polarity"],
                        "summary": analyzer.summarize_text(text, num_sentences=2)
                    })
            except Exception as e:
                print(f"Error processing {file.filename}: {e}")
                results.append({
                    "filename": file.filename,
                    "error": str(e)
                })

    return jsonify({
        "total_files": len(files),
        "processed_files": len([r for r in results if "error" not in r]),
        "results": results
    })

@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    return jsonify({
        "status": "healthy",
        "timestamp": datetime.datetime.now().isoformat(),
        "version": "2.0.0"
    })

# --- Main execution ---
if __name__ == '__main__':
    print("Starting Enhanced DocuMentor Pro API...")
    print("Available endpoints:")
    print("  POST /analyze - Analyze a single document")
    print("  POST /batch-analyze - Analyze multiple documents")
    print("  GET /health - Health check")
    app.run(debug=True, port=5001)