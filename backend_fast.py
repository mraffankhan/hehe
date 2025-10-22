import io
import json
import base64
import datetime
from collections import Counter
import heapq

# Flask for web server
from flask import Flask, request, jsonify
from flask_cors import CORS

# Libraries for analysis
import PyPDF2
from pptx import Presentation
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords
from nltk.sentiment.vader import SentimentIntensityAnalyzer
import textstat
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.cluster import KMeans
import pandas as pd
import plotly.graph_objects as go
from wordcloud import WordCloud
import matplotlib.pyplot as plt
import numpy as np
from textblob import TextBlob

# Ensure NLTK data is available
try:
    nltk.data.find("tokenizers/punkt")
except LookupError:
    nltk.download("punkt")
try:
    nltk.data.find("corpora/stopwords")
except LookupError:
    nltk.download("stopwords")
try:
    nltk.data.find("sentiment/vader_lexicon")
except LookupError:
    nltk.download("vader_lexicon")

plt.switch_backend('Agg')

app = Flask(__name__)
CORS(app)

class FastDocuMentorProAnalyzer:
    def __init__(self):
        self.stop_words = set(stopwords.words('english'))
        self.sia = SentimentIntensityAnalyzer()
        # No LanguageTool for faster processing

    def extract_text(self, file):
        # ... keep the same extraction code from before ...
        text = ""
        filename = file.filename.lower()
        file_size = len(file.read())
        file.seek(0)
        
        try:
            if filename.endswith(".pdf"):
                pdf_reader = PyPDF2.PdfReader(file)
                num_pages = len(pdf_reader.pages)
                for page in pdf_reader.pages:
                    text += (page.extract_text() or "")
                metadata = {"pages": num_pages, "file_size": file_size, "file_type": "PDF"}
            elif filename.endswith(".pptx"):
                prs = Presentation(file)
                num_slides = len(prs.slides)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text += shape.text + "\n"
                metadata = {"pages": num_slides, "file_size": file_size, "file_type": "PPTX"}
            elif filename.endswith(".txt"):
                raw = file.read()
                text = raw.decode('utf-8', errors='replace')
                metadata = {"pages": 1, "file_size": file_size, "file_type": "TXT"}
            else:
                return None, {"error": "Unsupported file format"}
                
            return text.strip(), metadata
        except Exception as e:
            print(f"Error extracting text: {e}")
            return None, {"error": str(e)}

    def get_enhanced_stats(self, text):
        if not text or len(text.strip()) == 0:
            return self._get_empty_stats()
            
        words = word_tokenize(text)
        sentences = sent_tokenize(text)
        paragraphs = [p for p in text.split('\n\n') if p.strip()]
        
        char_count = len(text)
        word_count = len(words)
        sentence_count = len(sentences)
        paragraph_count = len(paragraphs)
        
        try:
            flesch_ease = textstat.flesch_reading_ease(text)
            flesch_grade = textstat.flesch_kincaid_grade(text)
            avg_sentence_length = word_count / sentence_count if sentence_count > 0 else 0
            avg_word_length = char_count / word_count if word_count > 0 else 0
        except Exception:
            flesch_ease = flesch_grade = 0
            avg_sentence_length = avg_word_length = 0

        unique_words = set(word.lower() for word in words if word.isalpha())
        lexical_diversity = len(unique_words) / word_count if word_count > 0 else 0
        
        return {
            "basic": {
                "word_count": word_count, "sentence_count": sentence_count, 
                "paragraph_count": paragraph_count, "character_count": char_count,
                "avg_sentence_length": round(avg_sentence_length, 1), 
                "avg_word_length": round(avg_word_length, 1),
            },
            "readability": {
                "flesch_reading_ease": round(flesch_ease, 1),
                "flesch_kincaid_grade": round(flesch_grade, 1),
            },
            "vocabulary": {
                "unique_words": len(unique_words),
                "lexical_diversity": round(lexical_diversity, 3),
            }
        }

    def _get_empty_stats(self):
        return {
            "basic": {"word_count": 0, "sentence_count": 0, "paragraph_count": 0, 
                     "character_count": 0, "avg_sentence_length": 0, "avg_word_length": 0},
            "readability": {"flesch_reading_ease": 0, "flesch_kincaid_grade": 0},
            "vocabulary": {"unique_words": 0, "lexical_diversity": 0}
        }

    def get_enhanced_sentiment(self, text):
        if not text:
            return self._get_empty_sentiment()
            
        vader_scores = self.sia.polarity_scores(text)
        blob = TextBlob(text)
        
        return {
            "vader_scores": {
                "positive": round(vader_scores['pos'] * 100, 2),
                "neutral": round(vader_scores['neu'] * 100, 2),
                "negative": round(vader_scores['neg'] * 100, 2),
                "compound": round(vader_scores['compound'], 3)
            },
            "textblob_scores": {
                "polarity": round(blob.sentiment.polarity, 3),
                "subjectivity": round(blob.sentiment.subjectivity, 3)
            }
        }

    def _get_empty_sentiment(self):
        return {
            "vader_scores": {"positive": 0, "neutral": 100, "negative": 0, "compound": 0},
            "textblob_scores": {"polarity": 0, "subjectivity": 0}
        }

    def get_writing_quality(self, text):
        """Fast writing quality without grammar checking"""
        if not text:
            return {"grammar_errors": 0, "error_categories": {}, "quality_score": 100, "error_density": 0}
        
        word_count = len(word_tokenize(text))
        sentences = sent_tokenize(text)
        
        # Simple quality estimation
        if word_count < 50:
            quality_score = 60
        elif word_count > 2000:
            quality_score = 80
        else:
            quality_score = 85
            
        return {
            "grammar_errors": 0,
            "error_categories": {},
            "quality_score": quality_score,
            "error_density": 0
        }

    def get_content_structure(self, text):
        if not text:
            return {"sections": 0, "headings": [], "paragraph_distribution": []}
            
        paragraphs = [p for p in text.split('\n\n') if p.strip()]
        para_lengths = [len(word_tokenize(p)) for p in paragraphs]
        
        return {
            "sections": min(len(paragraphs), 5),
            "headings": [],
            "paragraph_distribution": {
                "short_paras": len([p for p in para_lengths if p < 50]),
                "medium_paras": len([p for p in para_lengths if 50 <= p <= 150]),
                "long_paras": len([p for p in para_lengths if p > 150])
            },
            "avg_paragraph_length": round(sum(para_lengths) / len(para_lengths), 1) if para_lengths else 0
        }

    def summarize_text(self, text, num_sentences=5):
        if not text:
            return ""
        sentences = sent_tokenize(text)
        if len(sentences) <= num_sentences:
            return text
            
        word_frequencies = Counter(w.lower() for w in word_tokenize(text) 
                                 if w.isalpha() and w.lower() not in self.stop_words)
        if not word_frequencies:
            return "Could not generate summary."
            
        max_freq = max(word_frequencies.values())
        for word in word_frequencies:
            word_frequencies[word] /= max_freq
            
        sentence_scores = {}
        for sent in sentences:
            for word in word_tokenize(sent.lower()):
                if word in word_frequencies:
                    sentence_scores[sent] = sentence_scores.get(sent, 0) + word_frequencies[word]
            sentence_scores[sent] = sentence_scores.get(sent, 0) / len(word_tokenize(sent))
            
        summary_sentences = heapq.nlargest(num_sentences, sentence_scores, key=sentence_scores.get)
        return ' '.join(summary_sentences)

    def get_enhanced_keywords(self, text, top_n=20):
        if not text:
            return []
            
        words = [w.lower() for w in word_tokenize(text) 
                if w.isalpha() and w.lower() not in self.stop_words and len(w) > 2]
        word_freq = Counter(words)
        return [(word, freq) for word, freq in word_freq.most_common(top_n)]

    def generate_insights(self, text, stats, sentiment, keywords):
        insights = []
        if not text:
            return insights
            
        word_count = stats["basic"]["word_count"]
        readability = stats["readability"]["flesch_reading_ease"]
        polarity = sentiment["textblob_scores"]["polarity"]
        
        if readability > 80:
            insights.append("Your document is very easy to read, suitable for a general audience.")
        elif readability > 60:
            insights.append("The document has good readability, appropriate for most readers.")
        elif readability > 30:
            insights.append("Consider simplifying some sentences to improve readability.")
        else:
            insights.append("The document may be difficult for many readers to understand.")
            
        if word_count < 300:
            insights.append("This is a concise document. Consider adding more detail if needed.")
        elif word_count > 2000:
            insights.append("This is a comprehensive document. Consider breaking it into sections.")
            
        if polarity > 0.3:
            insights.append("The document has a positive tone, which engages readers well.")
        elif polarity < -0.3:
            insights.append("The negative tone might affect how readers perceive the content.")
        else:
            insights.append("The document maintains a neutral, professional tone.")
            
        if keywords:
            top_keywords = [kw[0] for kw in keywords[:3]]
            insights.append(f"The main topics appear to be: {', '.join(top_keywords)}")
            
        return insights

    def create_enhanced_visualizations(self, text, keywords, stats, sentiment):
        visualizations = {}
        
        # Word Cloud
        if keywords:
            try:
                wc = WordCloud(width=800, height=400, background_color="rgba(255, 255, 255, 0)",
                             colormap='viridis', mode="RGBA").generate_from_frequencies(dict(keywords[:30]))
                img_buffer = io.BytesIO()
                wc.to_image().save(img_buffer, format='PNG')
                img_str = base64.b64encode(img_buffer.getvalue()).decode()
                visualizations["wordcloud"] = f"data:image/png;base64,{img_str}"
            except Exception as e:
                print(f"WordCloud error: {e}")
                visualizations["wordcloud"] = None
        
        return visualizations

analyzer = FastDocuMentorProAnalyzer()

@app.route('/analyze', methods=['POST'])
def analyze_document():
    if 'file' not in request.files:
        return jsonify({"error": "No file part in the request"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No file selected for uploading"}), 400

    try:
        text, metadata = analyzer.extract_text(file)
        if not text:
            return jsonify({"error": "Could not extract text from the document."}), 500

        # Fast analysis without heavy processing
        stats = analyzer.get_enhanced_stats(text)
        sentiment = analyzer.get_enhanced_sentiment(text)
        writing_quality = analyzer.get_writing_quality(text)
        content_structure = analyzer.get_content_structure(text)
        summary = analyzer.summarize_text(text)
        keywords = analyzer.get_enhanced_keywords(text)
        insights = analyzer.generate_insights(text, stats, sentiment, keywords)
        visualizations = analyzer.create_enhanced_visualizations(text, keywords, stats, sentiment)

        result = {
            "filename": file.filename,
            "metadata": metadata,
            "stats": stats,
            "sentiment": sentiment,
            "writing_quality": writing_quality,
            "content_structure": content_structure,
            "summary": summary,
            "keywords": keywords,
            "insights": insights,
            "visualizations": visualizations,
            "analysis_timestamp": datetime.datetime.now().isoformat()
        }
        
        return jsonify(result)

    except Exception as e:
        print(f"An error occurred during analysis: {e}")
        return jsonify({"error": "An internal error occurred during analysis."}), 500

@app.route('/health', methods=['GET'])
def health_check():
    return jsonify({"status": "healthy", "timestamp": datetime.datetime.now().isoformat()})

if __name__ == '__main__':
    print("Starting Fast DocuMentor Pro API...")
    app.run(debug=True, port=5001)